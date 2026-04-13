#!/usr/bin/env python3
import sys
from pathlib import Path

try:
    from pptx import Presentation
except Exception:
    print("python-pptx is required: pip install python-pptx", file=sys.stderr)
    raise


def extract_text(input_path: Path, output_path: Path) -> None:
    prs = Presentation(str(input_path))
    lines = []
    for i, slide in enumerate(prs.slides, start=1):
        lines.append(f"# Slide {i}")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                text = shape.text.strip()
                if text:
                    lines.append(text)
        lines.append("")

    output_path.parent.mkdir(parents=True, exist_ok=True)
    output_path.write_text("\n".join(lines), encoding="utf-8")


if __name__ == "__main__":
    if len(sys.argv) != 3:
        print("Usage: extract_ppt_text.py <input.pptx> <output.txt>", file=sys.stderr)
        sys.exit(1)
    extract_text(Path(sys.argv[1]), Path(sys.argv[2]))
